import sys, boto3, random, base64, os, secrets, time, datetime, json
from flask import jsonify  


import ast 
from pprint import pprint

from meta import BaseConfig   
s3_resource = BaseConfig.s3_resource  
S3_LOCATION = BaseConfig.S3_LOCATION
S3_BUCKET_NAME = BaseConfig.S3_BUCKET_NAME
SCHEMA = BaseConfig.SCHEMA
DESIGN = BaseConfig.DESIGN




class current_user ():
    username = 'Chris'

def get_all_values(nested_dictionary):
    detected = 0
    for key, value in nested_dictionary.items():        
        if type(value) is dict:
            print ('DICT FOUND', value)
            if get_all_values(value) != 0:
                detected += get_all_values(value)
        else:
            if value == None or value == "":
                print('CHECK', key, value)
                detected += 1
                
    return detected 



def createPPT():  
    #proj = request.form ['proj']
    proj = 'ND'
    #ansOBJ = request.form ['ansOBJ']
    ansOBJ = None

    ansDict = json.loads(ansOBJ)
    print(ansDict)

    if proj == 'ND':
        head = 'National Dish'
    if proj == 'CV':
        head = 'Cooking Video'
    
    if get_all_values(ansDict) != 0: 
        print('ERROR')
        return jsonify({'error' : 100})    

    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Food Culture English"
    subtitle.text = head + " Presentation by " + current_user.username



    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = head + ': ' + ansDict['Dish']
    
    tf = body_shape.text_frame
    tf.text = 'Reasons'
    for r in ansDict['Reasons']:
        p = tf.add_paragraph()
        p.text = ansDict['Reasons'][r]
        p.level = 1

    count = 1
    for part in ansDict['Parts']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Reason ' + str(count)
        
        tf = body_shape.text_frame
        tf.text = ansDict['Reasons'][part]        
        for r in ansDict['Parts'][part]['kw']:
            p = tf.add_paragraph()            
            p.text = ansDict['Parts'][part]['kw'][r]
            p.level = 1
        
        count +=1    

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Final Comment'
    tf = body_shape.text_frame
    tf.text = ansDict['Final']

    print('PROCESSING PPT') 
    filename = current_user.username + proj + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename) 
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename   
                         
    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)  
    
    return jsonify({'pptLink' : pptLink})


def createPPT_RR(): 
    
    #ansOBJ = request.form ['ansOBJ']
    ansDict = startDictGlobal('RR')
    #ansDict = json.loads(ansOBJ)
    print(ansDict)
        
    if get_all_values(ansDict) != 0: 
        print('ERROR')
        return jsonify({'error' : 100})    

    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Food Culture English: Restaurant Review"
    subtitle.text = " Presentation by " + current_user.username


    for entry in ansDict:
        if entry == 'Intro':
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = entry 
            tf = body_shape.text_frame
            tf.text = 'Restaurant Details' 

            p = tf.add_paragraph()
            p.text = ansDict[entry]['Name']
            p.level = 1
            p = tf.add_paragraph()
            p.text = ansDict[entry]['Style']
            p.level = 1
            p = tf.add_paragraph()
            p.text = ansDict[entry]['Location']
            p.level = 1
            p = tf.add_paragraph()
            p.text = 'when I go/went'
            p.level = 1

            pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )
            
        else:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = entry 

            tf = body_shape.text_frame 
            tf.text = 'Key Words:'
            for word in ansDict[entry]['key words'].split('/'):                 
                p = tf.add_paragraph()
                p.text = word
                p.level = 1

            pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )

    

    print('PROCESSING PPT') 
    filename = current_user.username + 'RR' + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename) 
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename   
                         
    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)  
    
    return jsonify({'pptLink' : pptLink})


def createPPT_VV(): 
    
    startDictGlobal = {
            'Video Title' : 'vv_title', 
            'Video Length' : '199', 
            'Video Views' : '2 million', 
            'Video Link' : 'youtube.com',
            'Video Point' : 'my point',            
            'Warm Up Question' : 'my question is ......',
            'Your Answer' : 'my answer is ....',
            'Description' : 'None',
            'Description Key Words' : 'test1 / test2 / test3',  
            'Comments' : 'commentary', 
            'Comments Key Words' : 'comm1 / comm2 / comm3'
        }
    #ansOBJ = request.form ['ansOBJ']
    ansDict = startDictGlobal
    #ansDict = json.loads(ansOBJ)
    print(ansDict)
        
    if get_all_values(ansDict) != 0: 
        print('ERROR')
        return jsonify({'error' : 100})    

    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentation English: Viral Video"
    subtitle.text = " Presentation by " + current_user.username

        #intro        
    bullet_slide_layout = prs.slide_layouts[3]
    # layout [3] is the two boxes next to each other
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'A Viral Video'
    body_shape = shapes.placeholders[1]    
    tf = body_shape.text_frame 
    
    p = tf.add_paragraph()
    p.text = ansDict['Video Title']    
    p = tf.add_paragraph()
    p.text = ansDict['Video Length']
    p = tf.add_paragraph()
    p.text = ansDict['Video Views']
    pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )
    
    #question        
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Question'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame 
    p = tf.add_paragraph()
    p.text = ansDict['Warm Up Question']        
    pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )

    #answer
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Answer'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame 
    p = tf.add_paragraph()
    p.text = ansDict['Your Answer']        
    pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )
    
    #description
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Description'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame 
    tf.text = 'Key Words:'
    for word in ansDict['Description Key Words'].split('/'):                 
        p = tf.add_paragraph()
        p.text = word
        p.level = 1
    
    pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )
    
    #video
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Video'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame 
    tf.text = ansDict['Video Link'] 
    
    pf = shapes.add_picture('add_play.jpg', Inches(6), Inches(2) )
    
    #comments
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Comments'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame 
    tf.text = 'Key Words:'
    for word in ansDict['Comments Key Words'].split('/'):                 
        p = tf.add_paragraph()
        p.text = word
        p.level = 1
    
    pf = shapes.add_picture('add_image.png', Inches(6), Inches(2) )


    print('PROCESSING PPT') 
    filename = current_user.username + 'VV' + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename) 
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename   
                         
    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)  
    
    return jsonify({'pptLink' : pptLink})



def startDictGlobal(sd):

    midterm = ['ND', 'CV']

    if sd in midterm:
        startDictGlobal = {
            'Dish' : None,             
            'Link' : 'Video Link',             
            'Image' : 'Image Link',
            'Final' : None,
            'Writer' : None, 
            'Reasons' : {
                1 : None, 
                2 : None, 
                3 : None
            },
            'Parts' : { 
                1 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                2 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                3 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
            }
        }
    elif sd == 'RR': 
        startDictGlobal = {
            'Intro' : {
                'Name' : 'r_name', 
                'Style' : 'r_style', 
                'Location' : 'r_location', 
                'When' : 'every day'
            },
            'Menu' : {
                'Sentence 1' : 'test', 
                'Sentence 2' : 'test', 
                'Sentence 3' : 'test', 
                'key words' : 'test / test / test'
            },
            'Food' : {
                'Sentence 1' : 'test', 
                'Sentence 2' : 'test', 
                'Sentence 3' : 'test', 
                'key words' : 'test / test / test' 
            },
            'Decor' : {
                'Sentence 1' : 'test', 
                'Sentence 2' : 'test', 
                'Sentence 3' : 'test', 
                'key words' : 'test / test / test' 
            },
            'Atmosphere' : {
                'Sentence 1' : 'test', 
                'Sentence 2' : 'test', 
                'Sentence 3' : 'test', 
                'key words' : 'test / test / test'
            },
            'Rating' : {
                'Sentence 1' : 'test', 
                'Sentence 2' : 'test', 
                'Sentence 3' : 'test', 
                'key words' : 'test / test / test' 
            }   
        }
    else:
        startDictGlobal = {
            'Video Title' : None, 
            'Video Length' : None, 
            'Video Views' : None, 
            'Video Link' : None,
            'Video Point' : None,            
            'Warm Up Question' : None,
            'Your Answer' : None,
            'Description' : None,
            'Description Key Words' : None,  
            'Comments' : None, 
            'Comments Key Words' : None
        }
    
    return startDictGlobal


answer = createPPT_VV()




