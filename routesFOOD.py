import sys, boto3, random, base64, os, secrets, time, datetime, json
from sqlalchemy import asc, desc, func, or_
from flask import render_template, url_for, flash, redirect, request, abort, jsonify
from app import app, db, bcrypt, mail
from flask_login import login_user, current_user, logout_user, login_required
from forms import *
from models import *
import ast
from pprint import pprint
from routesUser import get_grades, get_sources

from meta import BaseConfig
s3_resource = BaseConfig.s3_resource
S3_LOCATION = BaseConfig.S3_LOCATION
S3_BUCKET_NAME = BaseConfig.S3_BUCKET_NAME
SCHEMA = BaseConfig.SCHEMA
DESIGN = BaseConfig.DESIGN


def get_food_projects():
    content_object = s3_resource.Object( S3_BUCKET_NAME, 'json_files/sources.json' )
    file_content = content_object.get()['Body'].read().decode('utf-8')
    sDict = json.loads(file_content)  # json loads returns a dictionary
    ##print(sDict)
    return sDict


@app.route ("/food_list", methods=['GET','POST'])
@login_required
def food_list():

    sDict = get_food_projects()
    source = sDict['1']['M2']

    link = ''
    grade = 0

    project_answers = U021U.query.filter_by(username=current_user.username).first()
    if project_answers:
        link = project_answers.Ans03
        grade = project_answers.Grade



    return render_template('food/food_list.html', legend='Food Projects', source=source, link=link, grade=grade)


def getShareList():
    projects = U021U.query.all()


    listps = []
    allps = {}

    exNames = ['Chris', 'Janice', 'Frank', 'Rebbeca', 'Kanny', 'Tim', 'Andy', 'Mandy', 'Yuling ', 'Ivan', 'Katherine', current_user.username]

    for p in projects:
        if p.Grade >= 6 and p.username not in exNames:
            allps[p.username] = p.Ans03
            listps.append(p.username)

    print(allps)

    random.shuffle(listps)

    shareps = {}

    count = 0

    for k in listps:
        if count < 3:
            shareps[k] = {}
            shareps[k]['link'] = allps[k]
            shareps[k]['answers'] = {
                'Restaurant info' : '',
                'Menu sentence 1' : '',
                'Menu sentence 2' : '',
                'Food sentence 1' : '',
                'Food sentence 2' : '',
                'Decor sentence 1' : '',
                'Decor sentence 2' : '',
                'Atmosphere sentence 1' : '',
                'Atmosphere sentence 2' : '',
                'Final comment' : '',
            }
            count += 1

    return json.dumps(shareps)


@app.route ("/food/sharing", methods=['GET','POST'])
@login_required
def food_sharing():

    project_answers = U021U.query.filter_by(username=current_user.username).first()


    if project_answers.Ans05 == None or project_answers.Ans05 == '{}' :
        project_answers.Ans05 = getShareList()
        db.session.commit()

    share_projects = project_answers.Ans05


    return render_template('food/share_list.html', legend='Restaurant Sharing', share_projects=share_projects)


@app.route('/updateShare', methods=['POST'])
def updateShare():
    print('TEST')
    obj = request.form ['shareOBJ']
    ug = request.form ['updateGrade']
    print(obj, ug)

    project_answers = U021U.query.filter_by(username=current_user.username).first()
    project_answers.Ans05 = obj
    if ug:
        project_answers.Grade = 7
        print('Grade updated')


    db.session.commit()

    return jsonify({'result' : True})


def get_all_values(nested_dictionary):
    detected = 0
    for key, value in nested_dictionary.items():
        if type(value) is dict:
            print ('DICT FOUND', value)
            if get_all_values(value) != 0:
                detected += get_all_values(value)
        else:
            if value == None or value == "":
                print('CHECK', key, value)
                detected += 1

    return detected

@app.route('/updateFood', methods=['POST'])
def updateFood():
    proj = request.form ['proj']
    ansOBJ = request.form ['ansOBJ']

    ansDict = json.loads(ansOBJ)
    print(proj, ansDict)
    try:
        ## instructor update
        grade = request.form ['grade']
        name = request.form ['name']
    except:
        name = current_user.username
        if get_all_values(ansDict) == 0:
            grade = 4
        else:
            print ('GET_ALL', get_all_values(ansDict))
            grade = 0

    print('GRADE', grade)


    if proj == 'ND':
        project_answers = U011U.query.filter_by(username=name).first()
        project_answers.Ans01 = json.dumps(ansDict)
    if proj == 'CV':
        project_answers = U011U.query.filter_by(username=name).first()
        project_answers.Ans02 = json.dumps(ansDict)
    if proj == 'RR':
        project_answers = U021U.query.filter_by(username=name).first()
        project_answers.Ans01 = json.dumps(ansDict)

    project_answers.Grade = grade
    db.session.commit()

    return jsonify({'grade' : grade})

@app.route('/updateLink', methods=['POST'])
def updateLink():
    link = request.form ['link']



    project_answers = U021U.query.filter_by(username=current_user.username).first()
    project_answers.Ans03 = link

    db.session.commit()

    return jsonify({'link' : link})

@app.route('/createPPT', methods=['POST'])
def createPPT():
    proj = request.form ['proj']
    ansOBJ = request.form ['ansOBJ']

    ansDict = json.loads(ansOBJ)
    print(ansDict)

    if proj == 'ND':
        head = 'National Dish'
    if proj == 'CV':
        head = 'Cooking Video'

    if get_all_values(ansDict) != 0:
        print('ERROR')
        return jsonify({'error' : 100})

    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Food Culture English"
    subtitle.text = head + " Presentation by " + current_user.username



    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = head + ': ' + ansDict['Dish']

    tf = body_shape.text_frame
    tf.text = 'Reasons'
    for r in ansDict['Reasons']:
        p = tf.add_paragraph()
        p.text = ansDict['Reasons'][r]
        p.level = 1

    count = 1
    for part in ansDict['Parts']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Reason ' + str(count)

        tf = body_shape.text_frame
        tf.text = ansDict['Reasons'][part]
        for r in ansDict['Parts'][part]['kw']:
            p = tf.add_paragraph()
            p.text = ansDict['Parts'][part]['kw'][r]
            p.level = 1

        count +=1

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Final Comment'
    tf = body_shape.text_frame
    tf.text = ansDict['Final']

    print('PROCESSING PPT')
    filename = current_user.username + proj + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename)
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename

    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)

    return jsonify({'pptLink' : pptLink})


@app.route('/createPPT_RR', methods=['POST'])
def createPPT_RR():

    ansOBJ = request.form ['ansOBJ']
    ansDict = json.loads(ansOBJ)
    print(ansDict)

    if get_all_values(ansDict) != 0:
        print('ERROR')
        return jsonify({'error' : 100})

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Food Culture English: Restaurant Review"
    subtitle.text = " Presentation by " + current_user.username


    for entry in ansDict:
        if entry == 'Intro':
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = entry
            tf = body_shape.text_frame
            tf.text = 'Restaurant Details'

            p = tf.add_paragraph()
            p.text = ansDict[entry]['Restaurant Name']
            p.level = 1
            p = tf.add_paragraph()
            p.text = ansDict[entry]['Style']
            p.level = 1
            p = tf.add_paragraph()
            p.text = ansDict[entry]['Location']
            p.level = 1
            p = tf.add_paragraph()
            p.text = 'when I go/went'
            p.level = 1


            pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )

        else:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = entry

            tf = body_shape.text_frame
            tf.text = 'Key Words:'
            for word in ansDict[entry]['key words'].split('/'):
                p = tf.add_paragraph()
                p.text = word
                p.level = 1

            pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )



    print('PROCESSING PPT')
    filename = current_user.username + 'RR' + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename)
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename

    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)

    return jsonify({'pptLink' : pptLink})





def startDictGlobal(sd):

    midterm = ['ND', 'CV']

    if sd in midterm:
        startDictGlobal = {
            'Dish' : None,
            'Link' : 'Video Link',
            'Image' : 'Image Link',
            'Final' : None,
            'Writer' : None,
            'Reasons' : {
                1 : None,
                2 : None,
                3 : None
            },
            'Parts' : {
                1 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                2 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                3 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
            }
        }
    else:
        startDictGlobal = {
            'Intro' : {
                'Restaurant Name' : None,
                'Style' : None,
                'Location' : None,
                'When' : None
            },
            'Menu' : {
                'Sentence 1' : None,
                'Sentence 2' : None,
                'Sentence 3' : None,
                'key words' : None
            },
            'Food' : {
                'Sentence 1' : None,
                'Sentence 2' : None,
                'Sentence 3' : None,
                'key words' : None
            },
            'Decor' : {
                'Sentence 1' : None,
                'Sentence 2' : None,
                'Sentence 3' : None,
                'key words' : None
            },
            'Atmosphere' : {
                'Sentence 1' : None,
                'Sentence 2' : None,
                'Sentence 3' : None,
                'key words' : None
            },
            'Rating' : {
                'Sentence 1' : None,
                'Sentence 2' : None,
                'Sentence 3' : None,
                'key words' : None
            }
        }

    return startDictGlobal

@app.route ("/food/<string:proj>", methods=['GET','POST'])
@login_required
def food_proj(proj):
    print(proj)
    sDict = get_food_projects()
    if proj == 'ND':
        title = 'National Dish'
        source = sDict['1']['M3']
        model = U011U
    elif proj == 'CV':
        source = sDict['1']['M4']
        title = 'Cooking Video'
        model = U011U
    else:
        source = sDict['1']['MA']
        title = 'Restaurant Review'
        model = U021U


    if model.query.filter_by(username=current_user.username).first():
        pass
    else:
        startDict = startDictGlobal(proj)
        start = model(username=current_user.username, Ans01=json.dumps(startDict), Ans02=json.dumps(startDict), Grade=0)
        db.session.add(start)
        db.session.commit()

    project = model.query.filter_by(username=current_user.username).first()
    if proj == 'ND':
        ansDict = project.Ans01
        html = 'food/food_proj_MT.html'
    if proj == 'CV':
        ansDict = project.Ans02
        html = 'food/food_proj_MT.html'
    if proj == 'RR':
        ansDict = project.Ans01
        html = 'food/food_proj_RR.html'


    return render_template(html, legend='Food Project',
    source=source, title=title, ansString=ansDict)



@app.route ("/food_MT", methods=['GET','POST'])
@login_required
def food_MT():

    users = User.query.all()

    mtDict = {}
    for user in users:
        mtDict[user.username] = {
            'ID' : user.studentID,
            'Dish' : None,
            'Proj' : None,
            'Grade' : 0,
            'Data' : startDictGlobal('ND')
        }

    project_answers = U011U.query.all()

    for answer in project_answers:
        proj = None
        data = {}
        dish = None
        ndDict = json.loads(answer.Ans01)
        cvDict = json.loads(answer.Ans02)
        print(get_all_values(ndDict))
        print(get_all_values(cvDict))

        if get_all_values(ndDict) < get_all_values(cvDict):
            dish = ndDict['Dish']
            data = ndDict
            proj = 'ND'
        elif get_all_values(ndDict) > get_all_values(cvDict):
            dish = cvDict['Dish']
            data = cvDict
            proj = 'CV'

        if dish and answer.Grade == 0:
            tGrade = 1
        else:
            tGrade = answer.Grade

        mtDict[answer.username]['Dish'] = dish
        mtDict[answer.username]['Proj'] = proj
        mtDict[answer.username]['Grade'] = tGrade
        mtDict[answer.username]['Data'] = data

    pprint(mtDict)

    sDict = get_food_projects()
    source1 = sDict['1']['M3']
    source2 = sDict['1']['M4']


    return render_template('food/food_MT.html', legend='Food Project',
    source1=source1, source2=source2, ansString=json.dumps(mtDict)  )



@app.route ("/food_FN", methods=['GET','POST'])
@login_required
def food_FN():

    users = User.query.all()

    fnDict = {}
    for user in users:
        fnDict[user.username] = {
            'ID' : user.studentID,
            'Grade' : 0,
            'Data' : startDictGlobal('RR'),
            'Link' : ''
        }

    project_answers = U021U.query.all()

    for answer in project_answers:
        rrDict = json.loads(answer.Ans01)
        print(get_all_values(rrDict))


        fnDict[answer.username]['Grade'] = answer.Grade
        fnDict[answer.username]['Data'] = rrDict
        fnDict[answer.username]['Link'] = answer.Ans03
        if answer.Ans05:
            fnDict[answer.username]['answers'] = answer.Ans05
        else:
            fnDict[answer.username]['answers'] = json.dumps({})

    pprint(fnDict)

    sDict = get_food_projects()
    source = sDict['1']['M2']


    return render_template('food/food_FN.html', legend='Food Project',
    source=source, ansString=json.dumps(fnDict)  )



