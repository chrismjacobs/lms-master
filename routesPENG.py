import sys, boto3, random, base64, os, secrets, time, datetime, json
from sqlalchemy import asc, desc, func, or_
from flask import render_template, url_for, flash, redirect, request, abort, jsonify
from app import app, db, bcrypt, mail
from flask_login import login_user, current_user, logout_user, login_required
from forms import *
from models import *
import ast
from pprint import pprint
from routesUser import get_grades, get_sources

from meta import BaseConfig
s3_client = BaseConfig.s3_client
s3_resource = BaseConfig.s3_resource
S3_LOCATION = BaseConfig.S3_LOCATION
S3_BUCKET_NAME = BaseConfig.S3_BUCKET_NAME
SCHEMA = BaseConfig.SCHEMA
DESIGN = BaseConfig.DESIGN


from routesFOOD import get_all_values

def get_peng_projects():
    content_object = s3_resource.Object( S3_BUCKET_NAME, 'json_files/sources.json' )
    file_content = content_object.get()['Body'].read().decode('utf-8')
    sDict = json.loads(file_content)  # json loads returns a dictionary
    #source = sDict['1']['M2']
    source = sDict['1']['M3']
    return source

startDict = {
            'Product' : None,
            'Brand' : None,
            'Why' : None,
            'Image' : None,
            'Intro' : {
                1 : None,
                2 : None,
                3 : None
                },
            'Parts' : {
                'Product' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Features' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Demo' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Extra' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Close' : {
                1 : None,
                2 : None,
                3 : None
                },
            },
            'Cues' : {
                'Product' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Features' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Demo' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Extra' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Close' : {
                1 : None,
                2 : None,
                3 : None
                },
            },
        }





@app.route ("/peng_list", methods=['GET','POST'])
@login_required
def peng_list():

    setup = 'MT'
    source = get_peng_projects()

    projectData = {
        'MT' : U011U,
        'FH' : U011U,
        'FN' : U021U
    }

    if projectData[setup].query.filter_by(username=current_user.username).first():
        pass
    elif setup == 'MT':
        highDict = {
            'ProductName' : None,
            'Features' : None,
            'Materials' : None,
            'Image' : None,
            'Parts' : {
                'Intro' : {
                1 : None,
                2 : None
                },
                'Product' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Features' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Extra' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Close' : {
                1 : None,
                2 : None
                },
            },
            'Cues' : {
                'Intro' : {
                1 : None,
                2 : None
                },
                'Product' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Features' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Extra' : {
                1 : None,
                2 : None,
                3 : None
                },
                'Close' : {
                1 : None,
                2 : None
                }
            },
        }
        start = U011U(username=current_user.username, Ans01=json.dumps(highDict), Grade=0, Comment='0')
        db.session.add(start)
        db.session.commit()
    elif setup == 'FN':
        startDict = {
            'Video Title' : None,
            'Video Length' : None,
            'Video Views' : None,
            'Video Link' : None,
            'Video Point' : None,
            'Warm Up Question' : None,
            'Your Answer' : None,
            'Description' : None,
            'Description Key Words' : None,
            'Comments' : None,
            'Comments Key Words' : None
        }

        start = U021U(username=current_user.username, Ans01=json.dumps(startDict), Grade=0, Comment='0')
        db.session.add(start)
        db.session.commit()


    project = projectData[setup].query.filter_by(username=current_user.username).first()
    ansDict = project.Ans01
    grade = project.Grade
    stage = project.Comment

    return render_template('peng/peng_list.html', legend='Presentation Projects', source=source, stage=stage, grade=grade)

@app.route ("/peng_dash/<string:MTFN>", methods=['GET','POST'])
@login_required
def peng_dash(MTFN):
    source = get_peng_projects()

    projDict = {}

    if MTFN == 'FS':
        midterms = U011U.query.all()
        html = 'peng/peng_dash_fs.html'
        for user in midterms:
            userDict = json.loads(user.Ans01)
            userDict['stage'] = user.Comment
            projDict[user.username] = userDict
    if MTFN == 'MT':
        midterms = U011U.query.all()
        html = 'peng/peng_dash.html'
        for user in midterms:
            userDict = json.loads(user.Ans01)
            userDict['stage'] = user.Comment
            projDict[user.username] = userDict
    if MTFN == 'FN':
        html = 'peng/peng_dash_vv.html'
        midterms = U021U.query.all()
        for user in midterms:
            userDict = json.loads(user.Ans01)
            userDict['stage'] = user.Comment
            projDict[user.username] = userDict

    ansString = json.dumps(projDict)


    return render_template(html, title=MTFN, ansString=ansString, source=source)



@app.route ("/peng/<string:MTFN>/<string:page_stage>", methods=['GET','POST'])
@login_required
def peng_proj(MTFN, page_stage):
    source = get_peng_projects()
    # midterm or final
    if MTFN == 'MT':
        project = U011U.query.filter_by(username=current_user.username).first()
        html = 'peng/peng_demo'
    elif MTFN == 'FN':
        project = U021U.query.filter_by(username=current_user.username).first()
        html = 'peng/peng_video'
    elif MTFN == 'FH':
        project = U011U.query.filter_by(username=current_user.username).first()
        html = 'peng/peng_fuhsin'
    ansDict = project.Ans01
    grade = project.Grade
    stage = project.Comment

    print(source)
    print(page_stage)

    print(page_stage)
    return render_template(html + page_stage + '.html', legend='Presentation Project', source=source, ansString=ansDict)


def get_all_values(nested_dictionary):
    detected = 0
    for key, value in nested_dictionary.items():
        if type(value) is dict:
            #print ('DICT FOUND', value)
            if get_all_values(value) != 0:
                detected += get_all_values(value)
        else:
            if value == None or value == "" :
                #print(key, value)
                detected += 1

    return detected

@app.route('/updatePENG', methods=['POST'])
def updatePENG():
    proj = request.form ['proj']
    ansOBJ = request.form ['ansOBJ']
    stage = request.form ['stage']
    try:
        image_b64 = request.form ['image_b64']
        audio_b64 = request.form ['audio_b64']
    except:
        image_b64 = None
        audio_b64 = None

    try:
        user = request.form ['user']
    except:
        user = current_user.username

    print (proj, stage, user, ansOBJ)

    ansDict = json.loads(ansOBJ)

    if image_b64:
        print('PROCESSING IMAGE')
        image = base64.b64decode(image_b64)
        filename = proj + '/' + user + '.png'
        imageLink = S3_LOCATION + filename
        s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=filename, Body=image)
        ansDict['Image'] = imageLink
    if audio_b64:
        print('PROCESSING AUDIO')
        audio = base64.b64decode(audio_b64)
        filename = proj + '/' + user + '.mp3'
        audioLink = S3_LOCATION + filename
        s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=filename, Body=audio)
        ansDict['Audio'+ stage] = audioLink

    if proj == 'MT' or proj == 'FH':
        project_answers = U011U.query.filter_by(username=user).first()
    if proj == 'FN':
        print('FN_done')
        print(ansOBJ)
        project_answers = U021U.query.filter_by(username=user).first()
    ansOBJ = json.dumps(ansDict)
    project_answers.Ans01 = str(ansOBJ)
    db.session.commit()


    ### check stage 1
    if int(stage) == 1:

        stage1_checklist = [
            ansDict['ProductName'],
            ansDict['Features'],
            ansDict['Materials'],
        ]

        '''
        stage1_checklist = [
            ansDict['Video Title'],
            ansDict['Video Length'],
            ansDict['Video Views'],
            ansDict['Video Link'],
            ansDict['Video Point'],
            ]
        '''


        if int(project_answers.Comment) > 1 :
            pass
        elif None in stage1_checklist:
            project_answers.Comment = 1  ## stage
            db.session.commit()
        else:
            project_answers.Comment = 2  ## stage
            db.session.commit()

    if int(stage) == 2:
        if get_all_values(ansDict) != 0:
            return jsonify({'fail' : True})
        else:
            print ('clear', get_all_values(ansDict))

        if int(project_answers.Comment) > 2 :
            pass
        else:
            project_answers.Comment = 3  ## stage
            db.session.commit()

    if int(stage) > 3:
        if int(project_answers.Comment) > int(stage):
            pass
        else:
            project_answers.Comment = stage  ## stage
            db.session.commit()

    print('stage', project_answers.Comment)
    return jsonify({'ansString' : str(ansOBJ), 'stage' : project_answers.Comment})


@app.route('/createPPT_VV', methods=['POST'])
def createPPT_VV():
    ansOBJ = request.form ['ansOBJ']
    ansDict = json.loads(ansOBJ)
    print(ansDict)

    if get_all_values(ansDict) != 0:
        print('ERROR')
        return jsonify({'error' : 100})

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentation English: Viral Video"
    subtitle.text = " Presentation by " + current_user.username

        #intro
    bullet_slide_layout = prs.slide_layouts[3]
    # layout [3] is the two boxes next to each other
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'A Viral Video'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = 'Title: ' + ansDict['Video Title']
    p = tf.add_paragraph()
    p.text = 'Time: ' + ansDict['Video Length']
    p = tf.add_paragraph()
    p.text = 'Views: ' + ansDict['Video Views']
    pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )

    #question
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Question'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = ansDict['Warm Up Question']
    pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )

    #answer
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Answer'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = ansDict['Your Answer']
    pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )

    #description
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Description'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = 'Key Words:'
    for word in ansDict['Description Key Words'].split('/'):
        p = tf.add_paragraph()
        p.text = word
        p.level = 1

    pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )

    #video
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Video'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    tf.text = 'Okay, now let`s enjoy the video before I give my comments'
    p.text = ansDict['Video Link']

    pf = shapes.add_picture('static/images/add_play.jpg', Inches(6), Inches(2) )

    #comments
    bullet_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'My Comments'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = 'Key Words:'
    for word in ansDict['Comments Key Words'].split('/'):
        p = tf.add_paragraph()
        p.text = word
        p.level = 1

    pf = shapes.add_picture('static/images/add_image.png', Inches(6), Inches(2) )


    print('PROCESSING PPT')
    filename = current_user.username + 'VV' + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename)
    data = open(filename, 'rb')
    aws_filename = 'MT/' + filename
    pptLink = S3_LOCATION + aws_filename

    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)

    return jsonify({'pptLink' : pptLink})
